import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches

# --- 解决中文显示问题 ---
plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei', 'FangSong', 'Arial Unicode MS', 'DejaVu Sans']  # 指定默认字体
plt.rcParams['axes.unicode_minus'] = False  # 解决负号 '-' 显示为方块的问题

# ----------------------------
# 1. 配置文件与时间点
# ----------------------------
file_config = [
    # {'date': '2025-05-26', 'file': 'zero.xlsx', 'name': 'V8004'},
    {'date': '2025-08-11', 'file': 'first.xlsx', 'name': 'V8101'},
    {'date': '2025-08-20', 'file': 'second.xlsx', 'name': 'V8102'},
    {'date': '2025-09-05', 'file': 'third.xlsx', 'name': 'V8104'},
]

output_dir = "interface_performance_analysis"
os.makedirs(output_dir, exist_ok=True)

# ----------------------------
# 2. 读取并处理每个文件
# ----------------------------
dfs = []

# 修改字段处理部分
for config in file_config:
    df = pd.read_excel(config['file'], sheet_name=0)
    df = df[df['模块'] == df['模块']]  # 过滤无效行
    # 保留 '模块' 列的值 **不在** 给定列表中的行
    df = df[~df['模块'].isin(['iapps', 'ipaas', 'ipaascatalog', 'ipluto', 'icluster', 'idevops', 'iearth', 'imesh', 'icloudhub', 'iops'])]
    df = df[df['次数'] > 200]

    # 重命名并保留 '请求方式'
    df = df.rename(columns={
        '请求路径': '请求路径',
        '次数': '请求次数',
        '平均值': 'AVG',
        '50分位': 'TP50',
        '95分位': 'TP95'
    })

    # 添加复合唯一键
    df['接口标识'] = df['请求方式'] + ' ' + df['请求路径']

    df['日期'] = pd.to_datetime(config['date'])
    df['时间标签'] = config['name']

    df = df[['接口标识', '请求方式', '请求路径', '请求次数', 'AVG', 'TP50', 'TP95', '日期', '时间标签']]
    dfs.append(df)
# ----------------------------
# 3. 合并所有数据
# ----------------------------
df_all = pd.concat(dfs, ignore_index=True)
df_all.sort_values(['接口标识', '日期'], inplace=True)

print("✅ 数据合并完成，共", len(df_all), "条接口性能记录")
print("📅 时间跨度：2025-03-01 至 2025-09-01")

df = pd.read_csv('api_logs.csv')  # 替换为实际路径
# 构建接口唯一标识
df['接口标识'] = df['request_method'] + " " + df['url_template']
df = df.rename(columns={
        'url_template': '请求路径',
        'request_method': '请求方式'
})
df = df[['请求方式', '请求路径', '接口标识']].drop_duplicates()

df_deduplicate = df_all[['请求方式', '请求路径', '接口标识']].drop_duplicates()
df_tmp = pd.concat([df_deduplicate, df], ignore_index=True)
df_deduplicate = df_tmp[['请求方式', '请求路径', '接口标识']].drop_duplicates()
df_deduplicate.to_excel("deduplicate.xlsx", sheet_name="sheet", index=False)

for index, df in enumerate(dfs):
    df_notouch = df_deduplicate[~df_deduplicate['接口标识'].isin(df['接口标识'])]
    df_notouch.to_excel(f"notouch_{index}.xlsx", sheet_name="sheet", index=False)

print("done")


# ----------------------------
# 4. 生成趋势图（关键指标）
# ----------------------------
def plot_trend_filtered(df, metric, title, filename, min_val=50, max_val=10000, top_n=50):
    """
    绘制过滤后的趋势图
    - min_val, max_val: 过滤 metric 的范围
    - top_n: 最多显示前 N 个接口
    """
    # 获取每个接口在“V8104”时间点的 metric 值
    latest = df[df['时间标签'] == 'V8104'][['接口标识', metric]].set_index('接口标识')[metric]

    # 筛选出 metric 在合理区间的接口
    valid_interfaces = latest[(latest >= min_val) & (latest <= max_val)].index

    df_filtered = df[df['接口标识'].isin(valid_interfaces)]

    # 只绘制有完整三个时间点的接口
    path_counts = df_filtered.groupby('接口标识').size()
    full_paths = path_counts[path_counts == 3].index
    df_filtered = df_filtered[df_filtered['接口标识'].isin(full_paths)]

    # 按V8104 metric 值排序，取 top_n
    top_interfaces = df_filtered[df_filtered['时间标签'] == 'V8104'] \
        .set_index('接口标识')[metric] \
        .sort_values(ascending=False) \
        .head(top_n).index

    df_plot = df_filtered[df_filtered['接口标识'].isin(top_interfaces)]

    # 开始绘图
    plt.figure(figsize=(30, 15), constrained_layout=True)

    # 在 plot_trend_filtered 内部，判断是否为“退步”
    df_compare1 = df_all.pivot(index='接口标识', columns='时间标签', values='TP95')
    df_compare1['变化'] = df_compare1['V8104'] - df_compare1['V8101']
    regressed = df_compare1[df_compare1['变化'] > 500].index  # 退步超过1秒

    # 绘图时：
    for interface in df_plot['接口标识'].unique():
        data = df_plot[df_plot['接口标识'] == interface]
        label = interface
        color = 'red' if interface in regressed else 'blue'
        ls = '-' if interface in regressed else '--'
        # print(data[['接口标识', '时间标签', '日期', metric]])
        # print(color)
        # print(ls)
        # print("=======================")
        plt.plot(data['日期'], data[metric], marker='o', label=label, color=color, linestyle=ls)

    plt.title(title, fontsize=14)
    plt.ylabel(f"{metric} (ms)")
    plt.xlabel("时间")
    plt.xticks(rotation=45, ha='right')

    # 图例放在右侧
    plt.legend(bbox_to_anchor=(1.05, 1), loc='best', fontsize='small')

    # 保存时确保完整
    plt.savefig(f"{output_dir}/{filename}.png", dpi=150, bbox_inches='tight', pad_inches=0.5)
    plt.close()


# 调用示例
plot_trend_filtered(df_all, 'AVG', '各接口AVG趋势（50-5000ms）', 'trend_avg_latency_filtered')
plot_trend_filtered(df_all, 'TP95', '各接口TP95延迟趋势（50-5000ms）', 'trend_tp95_filtered')


# ----------------------------
# 5. 计算 TP95 变化（V8104 vs V8101）
# ----------------------------
def comparison2csv (df, column):
    df_earliest = df[df['日期'] == '2025-08-11'][['接口标识', column]].rename(columns={column: column + '_V8101'})
    df_latest = df[df['日期'] == '2025-09-05'][['接口标识', column]].rename(columns={column: column + '_V8104'})

    df_compare = pd.merge(df_earliest, df_latest, on='接口标识')
    df_compare[column + '_变化'] = df_compare[column + '_V8104'] - df_compare[column + '_V8101']
    df_compare[column + '_变化率'] = (df_compare[column + '_变化'] / df_compare[column + '_V8101']) * 100
    df_compare = df_compare.sort_values(column + '_变化', ascending=False).round(2)

    # 保存完整对比
    df_compare.to_csv(f"{output_dir}/change_comparison_{column}.csv", index=False)

    # ----------------------------
    # 6. Top 10 退步 & 优化接口
    # ----------------------------
    top_deteriorate = df_compare.head(10)
    top_improve = df_compare.tail(10).sort_values(column + '_变化').round(2)

    top_deteriorate.to_csv(f"{output_dir}/top10_deteriorated_{column}.csv", index=False)
    top_improve.to_csv(f"{output_dir}/top10_improved_{column}.csv", index=False)
    return df_compare, top_deteriorate, top_improve


df_compare_p95, top_deteriorate_p95, top_improve_p95 = comparison2csv(df_all, 'TP95')
df_compare_avg, top_deteriorate_avg, top_improve_avg = comparison2csv(df_all, 'AVG')


def heatmap(df, column):
    # 使用 pivot_table 替代 pivot，避免重复索引问题
    heatmap_data = df.pivot_table(
        index='接口标识',
        columns='时间标签',
        values=column,
        aggfunc='mean'
    )
    # 可选：只显示部分接口（避免太长）
    top_20 = df[df['时间标签'] == 'V8104'].nlargest(20, column)['接口标识']
    heatmap_data = heatmap_data.loc[heatmap_data.index.isin(top_20)]

    # 绘图时，只显示路径部分
    heatmap_data.index = [idx.split(' ', 1)[1] for idx in heatmap_data.index]

    plt.figure(figsize=(8, 10))
    sns.heatmap(heatmap_data, annot=True, fmt=".0f", cmap="YlOrRd", cbar_kws={'label': 'TP95 (ms)'})
    plt.title("Top 20 接口热力图" + column)
    plt.ylabel("接口路径")
    plt.xlabel("时间点")
    plt.savefig(f"{output_dir}/heatmap_top20_{column}.png", dpi=150, bbox_inches='tight')
    plt.close()


heatmap(df_all, 'TP95')
heatmap(df_all, 'AVG')


def plot_boxplot(df, metric, title, filename, output_dir='.'):
    """
    绘制箱线图。
    """
    # 设置输出目录
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # 使用Seaborn绘制箱线图
    plt.figure(figsize=(16, 9))
    sns.boxplot(x='接口标识', y=metric, data=df)
    plt.title(title)
    plt.xticks(rotation=45, ha='right')  # 调整x轴标签以避免重叠
    plt.tight_layout()

    # 保存图像
    plt.savefig(f"{output_dir}/{filename}.png", bbox_inches='tight', dpi=300)
    plt.close()


def filter_data_for_boxplot(df, metric, min_requests=100):
    """
    过滤数据用于绘制箱线图。
    移除异常值，并根据最小请求量筛选接口。
    """
    # 使用IQR移除异常值
    # Q1 = df[metric].quantile(0.25)
    # Q3 = df[metric].quantile(0.75)
    # IQR = Q3 - Q1
    # df_filtered = df[~((df[metric] < (Q1 - 1.5 * IQR)) | (df[metric] > (Q3 + 1.5 * IQR)))]
    # 按接口请求量筛选
    return df[df['请求次数'] >= min_requests]


# 示例调用
# 示例调用
df_filtered_p95_for_boxplot = filter_data_for_boxplot(df_all, 'TP95', 5000)
plot_boxplot(df_filtered_p95_for_boxplot, 'TP95', '各接口TP95延迟分布', 'boxplot_tp95', output_dir=output_dir)
df_filtered_avg_for_boxplot = filter_data_for_boxplot(df_all, 'AVG', 5000)
plot_boxplot(df_filtered_avg_for_boxplot, 'AVG', '各接口AVG延迟分布', 'boxplot_avg', output_dir=output_dir)


def plot_scatterplot(df, x_metric, y_metric, title, filename, output_dir='.'):
    """
    绘制散点图。
    """
    # 设置输出目录
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # 创建散点图
    plt.figure(figsize=(16, 9))
    plt.scatter(df[x_metric], df[y_metric], alpha=0.5)
    plt.title(title)
    plt.xlabel(x_metric)
    plt.ylabel(y_metric)
    # 可选：添加趋势线
    from scipy.stats import linregress
    slope, intercept, r_value, p_value, std_err = linregress(df[x_metric], df[y_metric])
    plt.plot(df[x_metric], intercept + slope * df[x_metric], color='red')  # 添加回归线
    plt.tight_layout()
    # 保存图像
    plt.savefig(f"{output_dir}/{filename}.png", bbox_inches='tight', dpi=300)
    plt.close()


# 示例调用
plot_scatterplot(df_all, '请求次数', 'AVG', '请求次数与AVG关系', 'scatterplot_requests_avg_latency', output_dir=output_dir)


# ----------------------------
# 8. 生成分析报告
# ----------------------------
report = f"""
📈 接口性能趋势分析报告
生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
数据来源: first.xlsx, second.xlsx, third.xlsx
时间跨度: V8101 → 三个月前 → V8104

📊 总体情况:
- 共分析接口数: {df_all['接口标识'].nunique()}
- 完整趋势数据接口数TP95: {len(df_compare_p95)}
- 完整趋势数据接口数Avg: {len(df_compare_avg)}

🔝 Top 5 性能退步最严重接口（TP95增加最多）:
{top_deteriorate_p95.head(5)[['接口标识', 'TP95_V8101', 'TP95_V8104', 'TP95_变化', 'TP95_变化率']].to_string(index=False)}

✅ Top 5 性能优化最明显接口（TP95下降最多）:
{top_improve_p95.head(5)[['接口标识', 'TP95_V8101', 'TP95_V8104', 'TP95_变化', 'TP95_变化率']].to_string(index=False)}

🔝 Top 5 性能退步最严重接口（Avg增加最多）:
{top_deteriorate_avg.head(5)[['接口标识', 'AVG_V8101', 'AVG_V8104', 'AVG_变化', 'AVG_变化率']].to_string(index=False)}

✅ Top 5 性能优化最明显接口（Avg下降最多）:
{top_improve_avg.head(5)[['接口标识', 'AVG_V8101', 'AVG_V8104', 'AVG_变化', 'AVG_变化率']].to_string(index=False)}
"""

with open(f"{output_dir}/analysis_report.txt", "w", encoding="utf-8") as f:
    f.write(report)

print("✅ 所有分析完成！结果已保存至:", output_dir)


# ----------------------------
# 9. 创建 PPT 报告
# ----------------------------
def add_slide(prs, title_text, image_path=None, table_data=None):
    slide_layout = prs.slide_layouts[5]  # 使用空白布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

    if image_path:
        left = Inches(1)
        top = Inches(1.5)
        slide.shapes.add_picture(image_path, left, top, width=Inches(8))

    if table_data:
        rows = len(table_data) + 1
        cols = len(table_data[0])
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(2)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 设置表头
        for col_idx, header in enumerate(table_data[0]):
            table.cell(0, col_idx).text = header

        # 填充表格内容
        for row_idx, row in enumerate(table_data[1:], start=1):
            for col_idx, cell in enumerate(row):
                table.cell(row_idx, col_idx).text = str(cell)


# 创建 PPT
prs = Presentation()

# 封面页
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "接口性能趋势分析报告"
subtitle.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

# 添加趋势图
# ✅ 使用你实际生成的文件名
add_slide(prs, "各接口AVG趋势（50-5000ms）", image_path=f"{output_dir}/trend_avg_latency_filtered.png")
add_slide(prs, "各接口TP95延迟趋势（50-5000ms）", image_path=f"{output_dir}/trend_tp95_filtered.png")

add_slide(prs, "接口TP95热力图", image_path=f"{output_dir}/heatmap_top20_TP95.png")
add_slide(prs, "接口AVG热力图", image_path=f"{output_dir}/heatmap_top20_AVG.png")

add_slide(prs, "接口TP95箱线图", image_path=f"{output_dir}/boxplot_tp95.png")
add_slide(prs, "接口AVG箱线图", image_path=f"{output_dir}/boxplot_avg.png")

# 添加 Top 10 退步接口
top_deteriorate_p95_data = [["接口标识", "TP95_V8101", "TP95_V8104", "TP95_变化",
                         "TP95_变化率"]] + top_deteriorate_p95.values.tolist()
add_slide(prs, "Top 10 TP95性能退步接口", table_data=top_deteriorate_p95_data)
# 添加 Top 10 优化接口
top_improve_p95_data = [["接口标识", "TP95_V8101", "TP95_V8104", "TP95_变化", "TP95_变化率"]] + top_improve_p95.values.tolist()
add_slide(prs, "Top 10 TP95性能优化接口", table_data=top_improve_p95_data)

top_deteriorate_avg_data = [["接口标识", "AVG_V8101", "AVG_V8104", "AVG_变化",
                             "AVG_变化率"]] + top_deteriorate_avg.values.tolist()
add_slide(prs, "Top 10 AVG性能退步接口", table_data=top_deteriorate_avg_data)
# 添加 Top 10 优化接口
top_improve_avg_data = [["接口标识", "AVG_V8101", "AVG_V8104", "AVG_变化", "AVG_变化率"]] + top_improve_avg.values.tolist()
add_slide(prs, "Top 10 AVG性能优化接口", table_data=top_improve_avg_data)

# 保存 PPT
ppt_output_path = f"{output_dir}/performance_analysis_report.pptx"
prs.save(ppt_output_path)
print(f"📄 PPT 报告已生成: {ppt_output_path}")

# ----------------------------
# 10. 打包为 ZIP 文件
# ----------------------------
import shutil

zip_output_path = f"{output_dir}.zip"
shutil.make_archive(output_dir, 'zip', output_dir)
print(f"📦 ZIP 包已生成: {zip_output_path}")