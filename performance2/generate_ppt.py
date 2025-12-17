import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import glob


def create_powerpoint_from_charts():
    """
    Create a PowerPoint presentation from API performance charts
    """
    # Check if charts directory exists
    if not os.path.exists("api_performance_charts"):
        print("❌ Charts directory 'api_performance_charts' does not exist.")
        print("Please run the performance analysis script first to generate charts.")
        return
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "API Performance Analysis Report"
    subtitle.text = "Generated from Performance Data"
    
    # Get all chart images
    chart_images = glob.glob("api_performance_charts/*.png")
    
    if not chart_images:
        print("❌ No chart images found in 'api_performance_charts' directory.")
        print("Please run the performance analysis script first to generate charts.")
        return
    
    print(f"📊 Found {len(chart_images)} chart images to include in presentation")
    
    # Define slide layout for content slides
    content_slide_layout = prs.slide_layouts[5]  # Blank layout
    
    # Process each chart image
    for i, image_path in enumerate(chart_images):
        # Extract chart info from filename
        filename = os.path.basename(image_path)
        chart_title = filename.replace(".png", "").replace("_", " ").title()
        
        # Add slide
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = chart_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add chart image
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(11.33)
        height = Inches(5.5)
        
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
            print(f"✅ Added '{filename}' to slide {i+2}")
        except Exception as e:
            print(f"❌ Error adding image '{filename}': {str(e)}")
            
            # Add placeholder text if image fails to load
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"Failed to load image:\n{filename}\n\nError: {str(e)}"
    
    # Save presentation
    output_file = "api_performance_report.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint presentation saved as '{output_file}'")
    print(f"📊 Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_powerpoint_from_charts()